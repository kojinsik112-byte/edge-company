# -*- coding: utf-8 -*-
"""㈜엣지컴퍼니 지명원 PPTX 빌더 — ACRO 제품 제조·공급 + 스마트홈 개발 (입주주관사 제외)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY=RGBColor(0x12,0x34,0x7A); BLUE=RGBColor(0x1F,0x6F,0xEB); ICE=RGBColor(0xEE,0xF3,0xFB)
INK=RGBColor(0x1A,0x22,0x30); SUB=RGBColor(0x5B,0x66,0x78); GOLD=RGBColor(0xB8,0x8A,0x2E)
LINE=RGBColor(0xE3,0xE8,0xF0); WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x1F,0x9D,0x6B)
DASHBG=RGBColor(0xF3,0xF6,0xFC); FILLY=RGBColor(0xFF,0xF7,0xDA); FILLT=RGBColor(0x7A,0x5D,0x00)
HF="맑은 고딕"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=13.333,7.5

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shadow=False,round_=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    return shp

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    tf.vertical_anchor=anchor
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        for (t,sz,c,b) in para:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(sz); f.color.rgb=c; f.bold=b; f.name=HF
    return tb

def header(s,no,title,en,tag="주관 제외 · ACRO 제품·스마트홈"):
    rect(s,0,0,SW,0.16,fill=NAVY)
    txt(s,0.7,0.55,9,0.4,[[(no+"  ",16,BLUE,True),(title,26,INK,True)]])
    txt(s,0.72,1.12,9,0.3,[[(en,11,RGBColor(0xAE,0xB7,0xC6),True)]])
    rect(s,0.72,1.5,0.62,0.045,fill=BLUE)
    txt(s,SW-4.7,0.6,4.0,0.3,[[("㈜엣지컴퍼니  |  지명원",10,SUB,True)]],align=PP_ALIGN.RIGHT)

def imgbox(s,x,y,w,h,icon,label,desc,tag):
    b=rect(s,x,y,w,h,fill=DASHBG,line=RGBColor(0xB9,0xC6,0xDD),lw=1.5,round_=True)
    b.line.dash_style=None
    txt(s,x,y+h*0.16,w,h*0.7,[
        [(icon,20,RGBColor(0x7F,0x8B,0xA0),False)],
        [(label,11.5,RGBColor(0x48,0x56,0x6B),True)],
        [(desc,9,RGBColor(0x92,0xA0,0xB5),False)],
        [("▣ "+tag,8.5,BLUE,True)],
    ],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=1.05)

def kpi(s,x,y,w,n,l,nsize=26):
    rect(s,x,y,w,1.15,fill=WHITE,line=LINE,lw=1,round_=True)
    txt(s,x,y+0.16,w,0.5,[[(n,nsize,BLUE,True)]],align=PP_ALIGN.CENTER)
    txt(s,x,y+0.7,w,0.4,[[(l,10.5,SUB,True)]],align=PP_ALIGN.CENTER,sp=1.0)

def card(s,x,y,w,h,title,body,tcolor=INK):
    rect(s,x,y,w,h,fill=WHITE,line=LINE,lw=1,round_=True)
    txt(s,x+0.22,y+0.2,w-0.44,0.4,[[(title,13.5,tcolor,True)]])
    txt(s,x+0.22,y+0.62,w-0.44,h-0.8,[[(body,11,SUB,False)]],sp=1.12)

def note(s,x,y,w,text,h=0.62):
    rect(s,x,y,w,h,fill=ICE,round_=True)
    txt(s,x+0.18,y+0.12,w-0.36,h-0.24,[[(text,10.3,RGBColor(0x3F,0x4D,0x63),False)]],anchor=MSO_ANCHOR.MIDDLE,sp=1.1)

def set_cell(c,text,sz=11,color=INK,bold=False,fill=None,align=PP_ALIGN.LEFT):
    c.margin_left=Inches(0.08); c.margin_right=Inches(0.06)
    c.margin_top=Inches(0.03); c.margin_bottom=Inches(0.03)
    c.vertical_anchor=MSO_ANCHOR.MIDDLE
    if fill is not None: c.fill.solid(); c.fill.fore_color.rgb=fill
    else: c.fill.solid(); c.fill.fore_color.rgb=WHITE
    tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.color.rgb=color; r.font.bold=bold; r.font.name=HF

def rec_table(s,x,y,w,headers,rows,colw,rowh=0.34,hh=0.36):
    nrows=len(rows)+1; ncols=len(headers)
    gt=s.shapes.add_table(nrows,ncols,Inches(x),Inches(y),Inches(w),Inches(hh+rowh*len(rows))).table
    gt.first_row=False; gt.horz_banding=False
    for j,cw in enumerate(colw): gt.columns[j].width=Inches(cw)
    gt.rows[0].height=Inches(hh)
    for j,h in enumerate(headers):
        set_cell(gt.cell(0,j),h,sz=10.5,color=WHITE,bold=True,fill=NAVY,
                 align=PP_ALIGN.LEFT if j==1 else PP_ALIGN.CENTER)
    for i,row in enumerate(rows,1):
        gt.rows[i].height=Inches(rowh)
        rfill=RGBColor(0xFA,0xFC,0xFF) if i%2==0 else WHITE
        for j,val in enumerate(row):
            big = (j==3 and isinstance(val,str) and val.replace(',','').isdigit() and int(val.replace(',',''))>=1000)
            set_cell(gt.cell(i,j),val,sz=10.3,color=BLUE if big else INK,bold=big,fill=rfill,
                     align=PP_ALIGN.LEFT if j==1 else PP_ALIGN.CENTER)
    return gt

def totbar(s,x,y,w,items):
    n=len(items); gap=0.18; cw=(w-gap*(n-1))/n
    for i,(num,lab) in enumerate(items):
        cx=x+i*(cw+gap)
        rect(s,cx,y,cw,0.95,fill=NAVY,round_=True)
        txt(s,cx,y+0.14,cw,0.5,[[(num,22,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,cx,y+0.62,cw,0.3,[[(lab,9.5,RGBColor(0xCA,0xDC,0xFC),True)]],align=PP_ALIGN.CENTER)

# ============ S1 표지 ============
s=slide()
rect(s,0,0,SW,SH,fill=WHITE)
rect(s,0,0,SW,SH,fill=ICE)  # light bg
rect(s,0,0,0.35,SH,fill=NAVY)
txt(s,1.0,0.95,11,0.4,[[("EDGE COMPANY CO., LTD.",15,BLUE,True)]])
txt(s,1.02,1.45,11,0.3,[[("ACRO 스마트홈 제품 제조·공급  ·  스마트홈 시스템 개발  ·  실내/경관 조명",11.5,SUB,True)]])
txt(s,1.0,2.5,11,0.4,[[("指 名 願",15,SUB,True)]])
txt(s,0.95,2.95,11,1.4,[[("지  명  원",62,INK,True)]])
rect(s,1.0,4.5,1.0,0.07,fill=BLUE)
txt(s,1.0,4.8,11.5,0.8,[
    [("㈜엣지컴퍼니",24,INK,True)],
    [("ACRO 스마트홈 제품 제조·공급 · 스마트홈 시스템 자체 개발 · 실내/경관 조명",13,SUB,True)],
],sp=1.25)
rect(s,1.0,6.35,11.3,0.02,fill=LINE)
txt(s,1.0,6.5,8,0.7,[
    [("울산광역시 울주군 청량읍 상남1길 28, 2동  ·  전국 22개 지사",10.5,SUB,False)],
    [("대표이사 고진식   |   대표번호 1533-3210 · 010-3345-6107",10.5,SUB,False)],
],sp=1.3)
txt(s,9.0,6.5,3.3,0.7,[
    [("제출처  각종 건설사·시행사 귀중",10.5,INK,True)],
    [("2026. __. __",10.5,SUB,False)],
],align=PP_ALIGN.RIGHT,sp=1.3)

# ============ S2 개요 ============
s=slide(); header(s,"01","회사 개요","Company Overview")
txt(s,0.7,1.85,12,1.5,[
    [("㈜엣지컴퍼니는 자체 스마트홈 브랜드 ",12.5,RGBColor(0x3A,0x44,0x56),False),("ACRO(아크로)",12.5,INK,True),
     ("를 직접 제조·공급하고,",12.5,RGBColor(0x3A,0x44,0x56),False)],
    [("실링팬·전동커튼·스위치·조명 디밍을 ",12.5,RGBColor(0x3A,0x44,0x56),False),("하나의 앱으로 통합 제어하는 스마트홈 시스템을 국내 최초로 자체 개발",12.5,INK,True),
     ("한 회사입니다.",12.5,RGBColor(0x3A,0x44,0x56),False)],
    [("실내조명·실외 경관조명의 설계·시공·납품까지 일괄 수행하며, 공동구매를 통해 다수 입주 단지에 제품을 공급해 왔습니다.",12.5,RGBColor(0x3A,0x44,0x56),False)],
],sp=1.4)
kpi(s,0.7,3.7,2.85,"19,294","공동구매 공급 세대수")
kpi(s,3.75,3.7,2.85,"52억","매출액 (최근)")
kpi(s,6.8,3.7,2.85,"BB / A","신용 / 현금흐름 (NICE D&B)")
kpi(s,9.85,3.7,2.85,"22","전국 지사")
txt(s,0.7,5.15,9,0.35,[[("3대 핵심 역량",14,INK,True)]])
rect(s,0.7,5.5,0.22,0.32,fill=BLUE)
card(s,0.7,5.95,3.9,1.2,"🌀 ACRO 제품 제조·공급","실링팬·조명·스위치·전동커튼·욕실 전열교환기를 자체 브랜드로 제조·공급.")
card(s,4.75,5.95,3.9,1.2,"📱 스마트홈 시스템 개발","팬·커튼·스위치·조명디밍 통합 앱 국내 최초 자체 개발.")
card(s,8.8,5.95,3.85,1.2,"💡 실내·경관 조명","주거·상업 실내조명 + 건물 외관 경관조명 설계·시공·납품.")

# ============ S3 일반현황 ============
s=slide(); header(s,"02","일반 현황","Company Profile")
info=[("상호","㈜엣지컴퍼니 (Edge Company Co., Ltd.)"),
("대표이사","고진식  (주주: 고진식 70% · 박은지 30%)"),
("설립일","2022년 2월 28일"),
("사업자등록번호","508-81-42798"),
("법인등록번호","230111-0369271"),
("본사 소재지","울산광역시 울주군 청량읍 상남1길 28, 2동  ·  전국 22개 지사"),
("자본금 / 상시인원","2억 원 (1주 10,000원 × 20,000주) / 14명"),
("주요 사업","ACRO 스마트홈 제품 제조·공급 · 스마트홈 시스템 개발 · 실내/경관 조명"),
("대표번호 / 대표이사","1533-3210  ·  010-3345-6107")]
ty=1.95
for k,v in info:
    rect(s,0.7,ty,3.0,0.5,fill=ICE,line=LINE,lw=0.75)
    txt(s,0.85,ty,2.8,0.5,[[(k,11,NAVY,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,3.7,ty,8.95,0.5,fill=WHITE,line=LINE,lw=0.75)
    txt(s,3.88,ty,8.7,0.5,[[(v,11,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    ty+=0.5
note(s,0.7,ty+0.12,11.95,"대표이사 고진식 — 온산 대한유화 10년 근무 후 창업, 저서 『사람 덕분에 여기까지 왔습니다』. 제조부터 현장까지 직접 챙기는 실무형 경영.",h=0.6)

# ============ S4 면허·인증 ============
s=slide(); header(s,"03","보유 면허 · 자격 · 인증","Licenses & Certifications")
txt(s,0.7,1.85,9,0.3,[[("면허 · 등록",13.5,INK,True)]]); rect(s,0.7,2.2,0.2,0.3,fill=BLUE)
lic=[("전기공사업 등록","제 울산-00821호 (2022.05.04 · 울산광역시장) — 조명 설치 전기결선 근거"),
("한국전기공사협회 회원","제00362호 (2022.05.09 · 정회원)"),
("건설업 / 소방시설업","[ 실제 등록번호 회장 확정 — KISCON 대조됨 ]")]
ty=2.65
for k,v in lic:
    rect(s,0.7,ty,3.2,0.5,fill=ICE,line=LINE,lw=0.75); txt(s,0.85,ty,3.0,0.5,[[(k,11,NAVY,True)]],anchor=MSO_ANCHOR.MIDDLE)
    fillv = FILLY if "회장 확정" in v else WHITE
    rect(s,3.9,ty,8.75,0.5,fill=fillv,line=LINE,lw=0.75)
    txt(s,4.08,ty,8.5,0.5,[[(v,10.5,FILLT if "회장 확정" in v else INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    ty+=0.5
txt(s,0.7,ty+0.18,9,0.3,[[("경영시스템 인증 · 신용",13.5,INK,True)]]); rect(s,0.7,ty+0.53,0.2,0.3,fill=BLUE)
ty2=ty+0.95
cert=[("ISO 9001:2015","품질경영시스템 (KS Q ISO 9001:2015)"),
("ISO 14001:2015","환경경영시스템 (KS I ISO 14001:2015)"),
("ISO 45001:2018","안전보건경영시스템 (KS Q ISO 45001:2018)"),
("신용평가","NICE디앤비 종합신용등급 BB · 현금흐름등급 A")]
for k,v in cert:
    rect(s,0.7,ty2,3.2,0.46,fill=ICE,line=LINE,lw=0.75); txt(s,0.85,ty2,3.0,0.46,[[(k,11,NAVY,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,3.9,ty2,8.75,0.46,fill=WHITE,line=LINE,lw=0.75); txt(s,4.08,ty2,8.5,0.46,[[(v,10.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    ty2+=0.46

# ============ S5 조직·인력 ============
s=slide(); header(s,"04","조직 & 인력","Organization")
imgbox(s,0.7,1.9,12.0,2.7,"🗂️","조직도",
       "대표이사 고진식 → [경영지원] [영업] [조명 시공본부] [ACRO 제품·생산본부] [스마트홈 개발] → 전국 22개 지사망",
       "ChatGPT/PPT 도형으로 생성")
txt(s,0.7,4.8,9,0.35,[[("기술 · 시공 인력 현황",14,INK,True)]]); rect(s,0.7,5.18,0.22,0.3,fill=BLUE)
labels=[("전기공사기술사","5명"),("조명·조도 설계","3명"),("시공팀·기능공","본사 8 · 지사 25"),("견적·적산","2명")]
cw=2.95; gap=0.13; x0=0.7
for i,(l,n) in enumerate(labels):
    cx=x0+i*(cw+gap)
    rect(s,cx,5.65,cw,1.2,fill=WHITE,line=LINE,lw=1,round_=True)
    txt(s,cx,5.82,cw,0.5,[[(n,17,BLUE,True)]],align=PP_ALIGN.CENTER)
    txt(s,cx,6.4,cw,0.35,[[(l,11,SUB,True)]],align=PP_ALIGN.CENTER)

# ============ S6 사업·제품 영역 ============
s=slide(); header(s,"05","사업 · 제품 영역","Business & Product Scope")
areas=[("🌀","ACRO 제품 제조·공급","실링팬·조명·스위치·전동커튼·욕실 전열교환기를 자체 브랜드로 제조·공급."),
("📱","스마트홈 시스템 개발","팬·커튼·스위치·조명 디밍을 하나의 앱으로 통합 제어 — 국내 최초 자체 개발."),
("💡","실내조명 시공·납품","주거·상업공간 조도·색온도 설계와 조명 설치·납품. 실링팬·디밍 조명 일괄."),
("🌃","실외 경관조명","건물 외관·단지 야간 경관조명 설계·시공. 단지 랜드마크화·야간 가치 향상.")]
positions=[(0.7,1.95),(6.75,1.95),(0.7,4.35),(6.75,4.35)]
for (ic,t,b),(x,y) in zip(areas,positions):
    rect(s,x,y,5.85,2.25,fill=WHITE,line=LINE,lw=1,round_=True)
    rect(s,x+0.3,y+0.3,0.75,0.75,fill=ICE,round_=True)
    txt(s,x+0.3,y+0.3,0.75,0.75,[[(ic,24,BLUE,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+1.25,y+0.4,4.4,0.5,[[(t,15,INK,True)]])
    txt(s,x+1.25,y+1.0,4.4,1.1,[[(b,11.5,SUB,False)]],sp=1.2)

# ============ S7 연혁 요약(24·25·26) ============
s=slide(); header(s,"06","공동구매 공급 연혁 — 요약","Group-Purchase Supply Record")
txt(s,0.7,1.9,12,0.8,[[("입주 단지 공동구매를 통해 ACRO·조명 제품을 공급한 연혁입니다. ",12.5,RGBColor(0x3A,0x44,0x56),False),
                       ("사송신도시 90% · 에코델타시티 90% 수임.",12.5,INK,True)]],sp=1.3)
yr=[("2024년","[ 단지 목록 회장 확정 ]","—","—"),
    ("2025년","10개 단지","10,146","4"),
    ("2026년 (확정)","12개 단지","9,148","2")]
rec_table(s,0.7,2.7,12.0,["연도","단지 수","총 세대수","1,000세대↑"],yr,[2.6,4.6,2.6,2.2],rowh=0.52,hh=0.42)
totbar(s,0.7,5.2,12.0,[("22","'25~'26 누적 단지"),("19,294","누적 세대수"),("90%","사송·에코델타 수임률"),("6","1,000세대↑ 대단지")])
note(s,0.7,6.45,12.0,"※ 2024년 단지 목록은 회장님 자료 입력 후 채웁니다. 2025·2026은 실제 수임실적 기준.",h=0.55)

# ============ S8 연혁 2025 ============
s=slide(); header(s,"06","공동구매 공급 연혁 ① 2025년","Record 2025")
rows25=[["1","울산 유보라 신천매곡","울산","352","공동구매"],
["2","춘천 중해마루 ★","강원 춘천","1,114","공동구매"],
["3","양정자이 SK뷰 ★","부산진구 양정동","2,276","공동구매"],
["4","에코델타시티 강서자이","부산 강서구","856","공동구매"],
["5","기장 유림노르웨이숲","부산 기장","165","공동구매"],
["6","두산위브 오션시티 ★","부산 남구 우암동","2,205","공동구매"],
["7","기장역 엘쿠르","부산 기장","219","공동구매"],
["8","에코델타시티 이편한세상","부산 강서구","953","공동구매"],
["9","에코델타 푸르지오린","부산 강서구","886","공동구매"],
["10","에코델타 대성베르힐 ★","부산 강서구","1,120","공동구매"]]
rec_table(s,0.7,1.85,12.0,["NO","단지명","지역","세대수","구분"],rows25,[0.9,4.6,3.3,1.7,1.5],rowh=0.345,hh=0.36)
totbar(s,0.7,5.75,12.0,[("10","공급 단지"),("10,146","총 세대수"),("4","1,000세대↑ 대단지")])

# ============ S9 연혁 2026 ============
s=slide(); header(s,"06","공동구매 공급 연혁 ② 2026년","Record 2026")
rows26=[["1","사송 데시앙 3차","양산 사송","672","확정"],
["2","사송 LH 신혼희망타운 2차","양산 사송","479","확정"],
["3","사송 더샵 3차","양산 사송","533","확정"],
["4","사송 제일 풍경채","양산 사송","452","확정"],
["5","사송 우미린","양산 사송","688","확정"],
["6","사송 롯데캐슬","양산 사송","—","진행"],
["7","에코델타 이편한세상 센터포인트","부산 강서구","953","확정"],
["8","에코델타 강서자이","부산 강서구","856","확정"],
["9","에코델타 푸르지오린","부산 강서구","886","확정"],
["10","에코델타 대성베르힐 ★","부산 강서구","1,120","확정"],
["11","에코델타 대방 디에트르(더퍼스트)","부산 강서구","972","확정"],
["12","에코델타 중흥S클래스(16블록) ★","부산 강서구","1,067","진행"],
["13","에코델타 대방 엘리움","부산 강서구","470","진행"]]
rec_table(s,0.7,1.8,12.0,["NO","단지명","지역","세대수","구분"],rows26,[0.9,5.0,2.9,1.6,1.5],rowh=0.295,hh=0.34)
totbar(s,0.7,6.1,12.0,[("22","'25~'26 누적"),("19,294","누적 세대수"),("90%","사송·에코델타 수임률")])

# ============ S10 ACRO 실링팬 ============
s=slide(); header(s,"07","ACRO 제품 ① 실링팬","Ceiling Fan")
def prod(s,x,y,w,h,name,tag,tagc,specs,price,old=None):
    rect(s,x,y,w,h,fill=WHITE,line=LINE,lw=1,round_=True)
    imgbox(s,x+0.18,y+0.17,3.0,h-0.34,"🌀",name.split()[-1],"회장 실제 사진","AI 생성 금지")
    tx=x+3.45
    txt(s,tx,y+0.22,w-3.7,0.4,[[(name,15,INK,True)]])
    txt(s,tx,y+0.66,w-3.7,0.3,[[(tag,10.5,tagc,True)]])
    yy=y+1.05
    for k,v in specs:
        txt(s,tx,yy,1.55,0.3,[[(k,10,SUB,True)]])
        txt(s,tx+1.6,yy,w-3.7-1.6,0.3,[[(v,10,INK,False)]])
        yy+=0.29
    pr=[("정상 "+old+"  ",10.5,RGBColor(0xAA,0xB3,0xC2),False),(price,15,INK,True)] if old else [(price,15,INK,True)]
    txt(s,tx,yy+0.02,w-3.7,0.4,[pr])
prod(s,0.7,1.8,12.0,2.6,"ACRO 시그니처 실링팬","★ 프리미엄 · 블루투스 · 초슬림",GOLD,
     [("날개 지름","132cm (52인치)"),("모터/풍속","BLDC DC모터 · 6단"),("제어","블루투스 앱+리모컨 · 1~6시간 타이머"),("A/S","5년")],
     "309,000원")
prod(s,0.7,4.55,12.0,2.6,"ACRO 슬림 실링팬","저천장 특화 · 합리형",BLUE,
     [("몸통 높이","14.5cm급 초슬림 (저천장 밀착)"),("모터/풍속","BLDC DC모터 · 6단"),("색상","우드톤 4색 / 조명 없음"),("A/S","5년")],
     "159,000원",old="189,000원")

# ============ S11 ACRO 조명·스위치·커튼·환기 ============
s=slide(); header(s,"08","ACRO 제품 ② 조명·스위치·커튼·환기","Lighting · Switch · Curtain · ERV")
items=[("🔆","ACRO 조명 · 디밍","COB 고연색 조명 + 앱/스위치로 밝기·색온도 무단계 디밍. 우물천장 간접조명 등."),
("🎛️","ACRO 스위치","실링팬·조명·디밍을 통합 제어하는 스위치. 아크로로 통일되는 스마트홈 허브."),
("🪟","ACRO 전동커튼","모터·앱 연동 전동커튼. 조명·실링팬과 함께 앱 하나로 제어."),
("💨","ACRO 욕실 전열교환기","욕실 환기·제습·열회수(ERV) 시스템. 결로·곰팡이 저감, 앱 연동 제어.")]
positions=[(0.7,1.95),(6.75,1.95),(0.7,4.45),(6.75,4.45)]
for (ic,t,b),(x,y) in zip(items,positions):
    rect(s,x,y,5.85,2.3,fill=WHITE,line=LINE,lw=1,round_=True)
    imgbox(s,x+0.2,y+0.2,1.9,1.9,ic,t.split()[-1],"실제 사진 우선","제품/공간컷")
    txt(s,x+2.3,y+0.3,3.35,0.5,[[(t,13.5,INK,True)]])
    txt(s,x+2.3,y+0.85,3.35,1.0,[[(b,10.3,SUB,False)]],sp=1.15)
    txt(s,x+2.3,y+1.85,3.35,0.3,[[("가격/스펙 회장 확정",9.5,FILLT,True)]])
note(s,0.7,6.95,12.0,"조명 디밍·스위치·전동커튼·욕실 전열교환기 정확 모델·스펙·가격은 회장 확정 후 채웁니다.",h=0.42)

# ============ S12 스마트홈 앱 ============
s=slide(); header(s,"09","ACRO 스마트홈 — 국내 최초 통합 앱","Smart-Home System")
rect(s,0.7,1.9,12.0,1.25,fill=NAVY,round_=True)
txt(s,1.0,2.05,11.4,0.35,[[("KOREA FIRST · 자체 개발",11,RGBColor(0xCA,0xDC,0xFC),True)]])
txt(s,1.0,2.4,11.4,0.7,[[("실링팬 · 전동커튼 · 스위치 · 조명 디밍을 하나의 앱으로 통합 제어",17,WHITE,True)],
                        [("— 국내 최초 자체 개발 스마트홈 시스템 기술",13,RGBColor(0xCA,0xDC,0xFC),True)]],sp=1.15)
phs=[("📱","홈 대시보드","기기 한눈에 제어"),("🌀","팬 + 조명 디밍","풍속·밝기·색온도"),("🪟","커튼·환기","전동커튼·전열교환기")]
pw=2.6; gap=0.5; x0=2.55
for i,(ic,l,d) in enumerate(phs):
    cx=x0+i*(pw+gap)
    b=rect(s,cx,3.4,pw,2.5,fill=DASHBG,line=RGBColor(0xB9,0xC6,0xDD),lw=1.5,round_=True)
    txt(s,cx,4.0,pw,1.4,[[(ic,24,RGBColor(0x7F,0x8B,0xA0),False)],[(l,11.5,RGBColor(0x48,0x56,0x6B),True)],[(d,9.5,RGBColor(0x92,0xA0,0xB5),False)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=1.2)
note(s,0.7,6.1,12.0,"📸 위 3컷은 회장 실제 앱 화면 자리 (주관사PT/img/app1~3.png · Desktop\\아크로 사진 · ACRO_App_Prototype_KR.html 캡처).",h=0.5)

# ============ S13 강점·협력 ============
s=slide(); header(s,"10","왜 엣지컴퍼니인가","Why Edge Company")
adv=[("🌀 자체 제품 제조","실링팬·조명·스위치·커튼·환기를 자체 브랜드로 제조·공급. 마진·하자 일원화."),
("📱 국내 최초 스마트홈","팬·커튼·스위치·디밍 통합 앱 자체 개발. 경쟁사가 못 따라오는 기술."),
("💡 실내·경관 조명","주거·상업 실내조명 + 야간 경관조명 설계·시공·납품 일괄."),
("📦 공동구매 공급망","'25~'26 22단지·19,294세대 공급. 사송·에코델타 90% 수임."),
("📈 견실한 재무","매출 52억 · NICE D&B 신용 BB · 현금흐름 A · 자본금 2억."),
("🌐 전국 22개 지사","본사 울산 + 전국 22개 지사망으로 현장 대응·납기 안정.")]
pos=[(0.7,1.95),(4.75,1.95),(8.8,1.95),(0.7,3.75),(4.75,3.75),(8.8,3.75)]
for (t,b),(x,y) in zip(adv,pos):
    card(s,x,y,3.85,1.65,t,b)
rect(s,0.7,5.75,12.0,1.35,fill=ICE,round_=True)
txt(s,0.95,5.95,11.5,0.35,[[("건설사·시행사 협력 제안",13.5,NAVY,True)]])
txt(s,0.95,6.4,11.5,0.65,[
    [("✓ 단지 ACRO 스마트홈 빌트인/옵션 패키지   ✓ 실내·경관 조명 설계·시공·납품",11,RGBColor(0x33,0x40,0x4F),False)],
    [("✓ 공동구매 제품 공급   ✓ 준공 후 장기 A/S 협약",11,RGBColor(0x33,0x40,0x4F),False)],
],sp=1.25)

# ============ S14 연락처 ============
s=slide(); header(s,"11","연락처","Contact")
ct=[("대표이사","고진식"),("대표번호","1533-3210"),("대표이사 직통","010-3345-6107"),
("본사","울산광역시 울주군 청량읍 상남1길 28, 2동 (전국 22개 지사)"),("이메일/홈페이지","[ 회장 확정 ]")]
ty=2.0
for k,v in ct:
    rect(s,0.7,ty,3.0,0.52,fill=ICE,line=LINE,lw=0.75); txt(s,0.85,ty,2.8,0.52,[[(k,11,NAVY,True)]],anchor=MSO_ANCHOR.MIDDLE)
    fillv=FILLY if "회장 확정" in v else WHITE
    rect(s,3.7,ty,8.95,0.52,fill=fillv,line=LINE,lw=0.75); txt(s,3.88,ty,8.7,0.52,[[(v,11,FILLT if "회장 확정" in v else INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    ty+=0.52
rect(s,0.7,5.1,12.0,1.7,fill=NAVY,round_=True)
txt(s,0.7,5.4,12.0,0.4,[[("EDGE COMPANY CO., LTD.",13,RGBColor(0xCA,0xDC,0xFC),True)]],align=PP_ALIGN.CENTER)
txt(s,0.7,5.85,12.0,0.5,[[("㈜엣지컴퍼니",24,WHITE,True)]],align=PP_ALIGN.CENTER)
txt(s,0.7,6.45,12.0,0.3,[[("ACRO 스마트홈 제품 제조·공급 · 스마트홈 시스템 개발 · 실내/경관 조명",11.5,RGBColor(0xCA,0xDC,0xFC),False)]],align=PP_ALIGN.CENTER)

out=os.path.join(os.path.dirname(__file__),"엣지컴퍼니_지명원.pptx")
prs.save(out)
print("SAVED",out,"slides",len(prs.slides._sldIdLst))
